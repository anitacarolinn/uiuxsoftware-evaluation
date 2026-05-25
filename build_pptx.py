"""Generate Traditional Chinese PPTX (A2-B1 TOCFL level) from design tools evaluation.

Output: design-tools-evaluation-zh.pptx
Style: 16:9 widescreen, cream theme, minimal text, visual-heavy.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============= COLORS (matching HTML cream theme) =============
BG = RGBColor(0xF7, 0xF1, 0xE3)         # cream background
SURFACE = RGBColor(0xFD, 0xFA, 0xF2)    # light surface
SURFACE2 = RGBColor(0xF0, 0xE8, 0xD3)   # elevated
TEXT = RGBColor(0x2A, 0x25, 0x20)       # warm dark
TEXT_DIM = RGBColor(0x6B, 0x63, 0x58)   # dim text
BORDER = RGBColor(0xD9, 0xCF, 0xB8)     # cream border
ACCENT = RGBColor(0x4F, 0x46, 0xE5)     # indigo accent

# Tool colors
FIGMA = RGBColor(0x7C, 0x3A, 0xED)      # purple
PAPER = RGBColor(0xDC, 0x26, 0x26)      # red
STITCH = RGBColor(0x15, 0x80, 0x3D)     # green
AISTUDIO = RGBColor(0x1D, 0x4E, 0xD8)   # blue
CLAUDE = RGBColor(0xB4, 0x53, 0x09)     # amber

SUCCESS = RGBColor(0x16, 0xA3, 0x4A)    # green check
DANGER = RGBColor(0xDC, 0x26, 0x26)     # red x

FONT_TC = "Microsoft JhengHei"  # Traditional Chinese system font

# ============= SETUP =============
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]  # fully blank layout


# ============= HELPERS =============
def add_bg(slide, color=BG):
    """Fill slide background."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_TC):
    """Add a text box with given style."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_pill(slide, x, y, w, h, text, *, fill_color, text_color=RGBColor(0xFF, 0xFF, 0xFF),
             size=14, bold=True, align=PP_ALIGN.CENTER):
    """Rounded colored pill (e.g., tool badge)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Emu(80000); tf.margin_right = Emu(80000)
    tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_TC
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


def add_card(slide, x, y, w, h, *, fill=SURFACE, border=BORDER, accent_color=None):
    """Light card with optional top accent bar."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(0.75)
    if accent_color is not None:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()
    return card


def add_label(slide, x, y, text, *, color=ACCENT, size=11):
    """Small uppercase label (eyebrow)."""
    tb = slide.shapes.add_textbox(x, y, Inches(8), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_TC
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return tb


def add_h2(slide, x, y, text, *, size=36):
    """Section heading."""
    return add_text(slide, x, y, Inches(12), Inches(0.8), text,
                    size=size, bold=True, color=TEXT)


def add_bullet_block(slide, x, y, w, h, items, *, head=None, head_color=ACCENT,
                     bullet_color=TEXT, text_color=TEXT, size=13, line_size=15):
    """A card with header and bullet list."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(120000); tf.margin_right = Emu(120000)
    tf.margin_top = Emu(120000); tf.margin_bottom = Emu(120000)

    first = True
    if head:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = head
        r.font.name = FONT_TC
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = head_color
        first = False

    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(6)
        # bullet symbol
        b = p.add_run()
        b.text = "• "
        b.font.name = FONT_TC
        b.font.size = Pt(line_size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        # content
        c = p.add_run()
        c.text = item
        c.font.name = FONT_TC
        c.font.size = Pt(line_size)
        c.font.color.rgb = text_color
    return tb


# ============= SLIDE 1: COVER =============
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)

    add_label(s, Inches(1), Inches(1.2), "設計工具比較 2026", size=14)

    add_text(s, Inches(1), Inches(1.7), Inches(11), Inches(2.2),
             "5 個 AI 設計工具", size=72, bold=True, color=TEXT)
    add_text(s, Inches(1), Inches(2.9), Inches(11), Inches(1.2),
             "1 份誠實的比較", size=60, bold=True, color=ACCENT)

    add_text(s, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
             "Figma Make · Paper · Stitch AI · AI Studio · Claude Design",
             size=18, color=TEXT_DIM)

    # audience pills
    add_pill(s, Inches(1), Inches(5.2), Inches(2.2), Inches(0.5),
             "給 Vibe Coder", fill_color=SURFACE, text_color=TEXT, size=13, bold=False)
    add_pill(s, Inches(3.4), Inches(5.2), Inches(2.2), Inches(0.5),
             "給 美編", fill_color=SURFACE, text_color=TEXT, size=13, bold=False)

    add_text(s, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
             "2026-05-01  ·  共 13 頁  ·  按方向鍵切換頁面",
             size=11, color=TEXT_DIM)


# ============= SLIDE 2: 五個工具 =============
def slide_overview():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_label(s, Inches(0.7), Inches(0.5), "概觀")
    add_h2(s, Inches(0.7), Inches(0.85), "五個工具一覽")

    tools = [
        ("Figma Make", FIGMA, "AI 生成設計", "$16/月", "美編最愛"),
        ("Paper", PAPER, "程式碼為主", "$16-20/月", "Vibe Coder"),
        ("Stitch AI", STITCH, "Google 出的", "免費", "兩種人都行"),
        ("AI Studio", AISTUDIO, "完整網站", "免費", "Vibe Coder"),
        ("Claude Design", CLAUDE, "AI 設計空間", "Pro 內含", "兩種人都行"),
    ]
    n = len(tools)
    gap = Inches(0.18)
    total = Inches(12)
    cw = (total - gap * (n - 1)) / n
    cx0 = Inches(0.7)
    cy = Inches(2.0)
    ch = Inches(4.4)

    for i, (name, color, kind, price, who) in enumerate(tools):
        x = cx0 + (cw + gap) * i
        add_card(s, x, cy, cw, ch, accent_color=color)
        # tool badge
        add_pill(s, x + Inches(0.3), cy + Inches(0.4), cw - Inches(0.6), Inches(0.5),
                 name, fill_color=color, size=14)
        # kind
        add_text(s, x + Inches(0.2), cy + Inches(1.2), cw - Inches(0.4), Inches(0.4),
                 kind, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)
        # price (big)
        add_text(s, x + Inches(0.2), cy + Inches(1.9), cw - Inches(0.4), Inches(1.0),
                 price, size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        # who
        add_text(s, x + Inches(0.2), cy + Inches(3.4), cw - Inches(0.4), Inches(0.4),
                 "適合：", size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), cy + Inches(3.7), cw - Inches(0.4), Inches(0.5),
                 who, size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
             "都是網頁版  ·  都用 AI  ·  方向不一樣",
             size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ============= SLIDE 3: 費用比較 =============
def slide_cost():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_label(s, Inches(0.7), Inches(0.5), "價錢")
    add_h2(s, Inches(0.7), Inches(0.85), "費用比較")

    rows = [
        ("Figma Make", FIGMA, "3 個檔案", "$16/月", "$55-90/月"),
        ("Paper", PAPER, "100 次/週", "$16-20/月", "之後推出"),
        ("Stitch AI", STITCH, "350+200 次/月", "之後推出", "—"),
        ("AI Studio", AISTUDIO, "完全免費", "API 計費", "Vertex AI"),
        ("Claude Design", CLAUDE, "—", "Claude Pro 內含", "Enterprise 內含"),
    ]
    headers = ["", "免費版", "個人/團隊", "公司版"]

    table = s.shapes.add_table(len(rows) + 1, 4,
                               Inches(0.7), Inches(1.9),
                               Inches(12), Inches(4.5)).table
    # header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE2
        tf = cell.text_frame
        tf.text = h
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            for r in para.runs:
                r.font.name = FONT_TC
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.color.rgb = TEXT_DIM

    for ri, (name, color, free, pro, ent) in enumerate(rows, start=1):
        for ci, val in enumerate([name, free, pro, ent]):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE
            tf = cell.text_frame
            tf.text = val
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in para.runs:
                    r.font.name = FONT_TC
                    r.font.size = Pt(15)
                    r.font.bold = (ci == 0)
                    r.font.color.rgb = color if ci == 0 else TEXT

    # set column widths
    table.columns[0].width = Inches(2.5)
    for i in (1, 2, 3):
        table.columns[i].width = Inches(3.17)


# ============= SLIDE 4: 操作感受 =============
def slide_editor():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_label(s, Inches(0.7), Inches(0.5), "編輯器體驗")
    add_h2(s, Inches(0.7), Inches(0.85), "用起來感覺如何？")

    headers = ["功能", "Figma Make", "Paper", "Stitch", "AI Studio", "Claude"]
    head_colors = [TEXT_DIM, FIGMA, PAPER, STITCH, AISTUDIO, CLAUDE]

    rows = [
        ("即時協作", "只能分享", "基本", "沒有", "只能分享", "沒有"),
        ("互動預覽", "進階", "還沒有", "頁面切換", "可直接跑", "AI 自動做"),
        ("元件庫", "需匯出 Figma", "用 MCP", "貼網址抓樣式", "沒有", "讀你的程式碼"),
        ("程式碼匯出", "React", "HTML/CSS", "7 種框架", "完整網站", "交給 Claude Code"),
        ("其他匯出", ".fig + 程式碼", "HTML/CSS", "貼到 Figma", "雲端部署", "PDF/PPTX/Canva"),
    ]

    table = s.shapes.add_table(len(rows) + 1, len(headers),
                               Inches(0.5), Inches(1.85),
                               Inches(12.4), Inches(4.6)).table
    # header
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE2
        tf = cell.text_frame
        tf.text = h
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            for r in para.runs:
                r.font.name = FONT_TC
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = head_colors[ci]

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE
            tf = cell.text_frame
            tf.text = val
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in para.runs:
                    r.font.name = FONT_TC
                    r.font.size = Pt(11)
                    r.font.bold = (ci == 0)
                    r.font.color.rgb = TEXT

    table.columns[0].width = Inches(2.0)
    for i in range(1, 6):
        table.columns[i].width = Inches(2.08)

    add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.4),
             "備註：所有工具都能透過 AI 產生任何框架的程式碼，這欄只列「內建匯出」",
             size=11, color=TEXT_DIM, align=PP_ALIGN.LEFT)


# ============= SLIDE 5: AI 功能 =============
def slide_ai():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_label(s, Inches(0.7), Inches(0.5), "AI 功能")
    add_h2(s, Inches(0.7), Inches(0.85), "AI 功能比較")

    headers = ["功能", "Figma Make", "Paper", "Stitch", "AI Studio", "Claude"]
    head_colors = [TEXT_DIM, FIGMA, PAPER, STITCH, AISTUDIO, CLAUDE]

    rows = [
        ("文字 → 設計", "First Draft", "用 MCP", "核心 (5 頁)", "核心 (整個 app)", "核心 (對話)"),
        ("留言", "✗", "✗", "✗", "點擊留言/改", "留言 + 直接改"),
        ("MCP", "讀+寫", "24 個工具", "stitch-mcp", "Gemini MCP", "用 Claude Code"),
        ("自動上線", "✗", "✗", "✗", "Cloud Run", "用 Claude Code"),
    ]

    table = s.shapes.add_table(len(rows) + 1, len(headers),
                               Inches(0.5), Inches(1.85),
                               Inches(12.4), Inches(4.0)).table
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE2
        tf = cell.text_frame
        tf.text = h
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            for r in para.runs:
                r.font.name = FONT_TC
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = head_colors[ci]

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE
            tf = cell.text_frame
            tf.text = val
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in para.runs:
                    r.font.name = FONT_TC
                    r.font.size = Pt(12)
                    r.font.bold = (ci == 0)
                    r.font.color.rgb = TEXT

    table.columns[0].width = Inches(2.0)
    for i in range(1, 6):
        table.columns[i].width = Inches(2.08)


# ============= SLIDE 6: 測試方法 =============
def slide_method():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_label(s, Inches(0.7), Inches(0.5), "測試方法")
    add_h2(s, Inches(0.7), Inches(0.85), "一個指令，五個工具")

    add_text(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.5),
             "我用「同樣的指令」測試 5 個工具，這樣比較才公平。",
             size=18, color=TEXT_DIM)

    # input box
    add_card(s, Inches(0.7), Inches(2.9), Inches(3.0), Inches(2.4), accent_color=ACCENT)
    add_text(s, Inches(0.85), Inches(3.1), Inches(2.7), Inches(0.4),
             "輸入", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(0.85), Inches(3.5), Inches(2.7), Inches(1.2),
             "1 個詳細指令", size=32, bold=True, color=TEXT)
    add_text(s, Inches(0.85), Inches(4.55), Inches(2.7), Inches(0.6),
             "災難救助儀表板", size=14, color=TEXT_DIM)

    # arrow
    add_text(s, Inches(3.9), Inches(3.7), Inches(0.6), Inches(1.0),
             "→", size=48, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # 5 outputs
    tools = [
        ("Figma Make", FIGMA),
        ("Paper", PAPER),
        ("Stitch", STITCH),
        ("AI Studio", AISTUDIO),
        ("Claude", CLAUDE),
    ]
    out_x = Inches(4.7)
    out_w = Inches(1.6)
    out_gap = Inches(0.07)
    for i, (name, color) in enumerate(tools):
        x = out_x + (out_w + out_gap) * i
        add_card(s, x, Inches(2.9), out_w, Inches(2.4), accent_color=color)
        add_text(s, x, Inches(3.6), out_w, Inches(0.5), name,
                 size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(4.2), out_w, Inches(0.4), "結果",
                 size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.6),
             "同樣的目標 → 看每個工具怎麼做",
             size=18, bold=True, color=TEXT, align=PP_ALIGN.CENTER)


# ============= PER-TOOL DEEP DIVE =============
def slide_tool(name, color, tagline, pros, cons, verdict, link_label, link):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_label(s, Inches(0.7), Inches(0.5), "測試筆記")

    # Header: badge + title
    add_pill(s, Inches(0.7), Inches(0.95), Inches(2.4), Inches(0.55),
             name, fill_color=color, size=16)
    add_text(s, Inches(3.3), Inches(0.95), Inches(9), Inches(0.6),
             "我的測試筆記", size=30, bold=True, color=TEXT)

    # Tagline
    add_text(s, Inches(0.7), Inches(1.75), Inches(12), Inches(0.5),
             tagline, size=16, color=TEXT_DIM)

    # Pros card
    pros_x = Inches(0.7)
    pros_y = Inches(2.5)
    pros_w = Inches(6.0)
    pros_h = Inches(3.4)
    add_card(s, pros_x, pros_y, pros_w, pros_h)
    add_text(s, pros_x + Inches(0.3), pros_y + Inches(0.2), pros_w - Inches(0.6), Inches(0.4),
             "+ 優點", size=13, bold=True, color=SUCCESS)
    for i, p in enumerate(pros):
        y = pros_y + Inches(0.7 + 0.85 * i)
        add_text(s, pros_x + Inches(0.3), y, Inches(0.3), Inches(0.4),
                 "+", size=18, bold=True, color=SUCCESS)
        add_text(s, pros_x + Inches(0.65), y + Inches(0.02), pros_w - Inches(1), Inches(0.7),
                 p, size=14, color=TEXT)

    # Cons card
    cons_x = Inches(7.0)
    cons_y = Inches(2.5)
    cons_w = Inches(6.0)
    cons_h = Inches(3.4)
    add_card(s, cons_x, cons_y, cons_w, cons_h)
    add_text(s, cons_x + Inches(0.3), cons_y + Inches(0.2), cons_w - Inches(0.6), Inches(0.4),
             "− 缺點", size=13, bold=True, color=DANGER)
    for i, c in enumerate(cons):
        y = cons_y + Inches(0.7 + 0.85 * i)
        add_text(s, cons_x + Inches(0.3), y, Inches(0.3), Inches(0.4),
                 "−", size=18, bold=True, color=DANGER)
        add_text(s, cons_x + Inches(0.65), y + Inches(0.02), cons_w - Inches(1), Inches(0.7),
                 c, size=14, color=TEXT)

    # Verdict bar
    add_card(s, Inches(0.7), Inches(6.1), Inches(12.3), Inches(0.95),
             accent_color=ACCENT)
    add_text(s, Inches(0.95), Inches(6.25), Inches(8), Inches(0.4),
             f"結論：{verdict}", size=13, bold=True, color=TEXT)
    add_text(s, Inches(0.95), Inches(6.6), Inches(11.8), Inches(0.4),
             f"{link_label}：{link}", size=11, color=ACCENT)


# ============= SLIDE 12: FIGMA PIVOT =============
def slide_figma_pivot():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_label(s, Inches(0.7), Inches(0.5), "最新消息", color=FIGMA)

    add_pill(s, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.5),
             "Figma", fill_color=FIGMA, size=14)
    add_text(s, Inches(2.4), Inches(0.95), Inches(11), Inches(0.6),
             "Figma 大轉向：開放給 AI", size=30, bold=True, color=TEXT)

    add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.6),
             "Figma 在 3 月和 4 月發了兩篇文章，把 Figma 開放給 AI 程式助手用。",
             size=15, color=TEXT_DIM)

    # Two cards
    cw = Inches(6.0)
    ch = Inches(3.0)
    cy = Inches(2.7)

    # Card 1: March
    add_card(s, Inches(0.7), cy, cw, ch, accent_color=FIGMA)
    add_text(s, Inches(0.95), cy + Inches(0.2), cw - Inches(0.5), Inches(0.3),
             "3 月 24 日 — Figma 畫布", size=11, bold=True, color=FIGMA)
    add_text(s, Inches(0.95), cy + Inches(0.55), cw - Inches(0.5), Inches(0.5),
             "AI 可以直接畫 Figma", size=18, bold=True, color=TEXT)
    add_bullet_block(s, Inches(0.95), cy + Inches(1.1), cw - Inches(0.5), Inches(1.7), [
        "use_figma：AI 直接寫 Figma 檔案",
        "Skills：用 Markdown 教 AI 你的設計規則",
    ], size=12, line_size=12)
    add_text(s, Inches(0.95), cy + Inches(2.55), cw - Inches(0.5), Inches(0.4),
             "目標：執行 — AI 幫你做設計", size=11, color=TEXT_DIM)

    # Card 2: April
    add_card(s, Inches(7.0), cy, cw, ch, accent_color=FIGMA)
    add_text(s, Inches(7.25), cy + Inches(0.2), cw - Inches(0.5), Inches(0.3),
             "4 月 28 日 — FigJam 白板", size=11, bold=True, color=FIGMA)
    add_text(s, Inches(7.25), cy + Inches(0.55), cw - Inches(0.5), Inches(0.5),
             "AI 在 FigJam 上規劃", size=18, bold=True, color=TEXT)
    add_bullet_block(s, Inches(7.25), cy + Inches(1.1), cw - Inches(0.5), Inches(1.7), [
        "generate_diagram：產生架構圖",
        "get_figjam：把白板帶回程式碼裡",
    ], size=12, line_size=12)
    add_text(s, Inches(7.25), cy + Inches(2.55), cw - Inches(0.5), Inches(0.4),
             "目標：規劃 — 寫程式前先想清楚", size=11, color=TEXT_DIM)

    # Bottom callout
    add_card(s, Inches(0.7), Inches(6.0), Inches(12.3), Inches(1.0),
             accent_color=FIGMA)
    add_text(s, Inches(0.95), Inches(6.15), Inches(11.8), Inches(0.4),
             "代表什麼？", size=11, bold=True, color=FIGMA)
    add_text(s, Inches(0.95), Inches(6.5), Inches(11.8), Inches(0.5),
             "Figma 現在是「最強 AI 設計工具」 — 不只給美編用，也適合 Vibe Coder。Beta 期間免費。",
             size=13, color=TEXT)


# ============= SLIDE 13: FINAL =============
def slide_final():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_label(s, Inches(0.7), Inches(0.5), "總結")
    add_h2(s, Inches(0.7), Inches(0.85), "最後總結")

    summaries = [
        ("Figma Make", FIGMA, "美編首選，簡單網站一鍵發布"),
        ("Paper", PAPER, "畫面最漂亮，可以細修"),
        ("Stitch AI", STITCH, "免費，產出最快"),
        ("AI Studio", AISTUDIO, "做能用的網站，一鍵上線"),
        ("Claude Design", CLAUDE, "互動最完整，Claude Pro 內含"),
    ]
    n = len(summaries)
    gap = Inches(0.18)
    total = Inches(12)
    cw = (total - gap * (n - 1)) / n
    cx0 = Inches(0.7)
    cy = Inches(2.0)
    ch = Inches(2.3)

    for i, (name, color, text) in enumerate(summaries):
        x = cx0 + (cw + gap) * i
        add_card(s, x, cy, cw, ch, accent_color=color)
        add_text(s, x + Inches(0.2), cy + Inches(0.4), cw - Inches(0.4), Inches(0.5),
                 name, size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), cy + Inches(1.0), cw - Inches(0.4), Inches(1.2),
                 text, size=12, color=TEXT, align=PP_ALIGN.CENTER)

    # Best combo
    add_card(s, Inches(0.7), Inches(4.7), Inches(12), Inches(1.5),
             fill=SURFACE2, accent_color=ACCENT)
    add_text(s, Inches(0.7), Inches(4.95), Inches(12), Inches(0.4),
             "最佳組合", size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.7),
             "Claude Design / Stitch  →  Paper  →  Claude Code",
             size=22, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.4),
             "（生成）          （細修）            （上線）",
             size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
             "更新日期：2026-05-01",
             size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ============= BUILD =============
slide_cover()
slide_overview()
slide_cost()
slide_editor()
slide_ai()
slide_method()

# Per-tool deep dives
slide_tool(
    "Figma Make", FIGMA,
    "簡單網站一鍵發布。適合單頁網站。",
    [
        "可以直接從 Figma Make 發布網站，網址是 .figma.site",
        "可以連 GitHub、Supabase 資料庫，也能裝 npm 套件",
        "可以用自己的網域名稱（要去買網域）",
    ],
    [
        "畫面不能用滑鼠改，只能用指令或寫程式",
        "沒有手機/平板/電腦的預覽切換",
    ],
    "適合做簡單單頁網站。複雜的就不適合。",
    "示範網站", "stem-stand-02084533.figma.site",
)

slide_tool(
    "Paper", PAPER,
    "畫面最漂亮，可以用滑鼠改細節。我的常用工具。",
    [
        "用 Claude MCP 連接，先生成設計，再進 Paper 細修",
        "畫面品質很乾淨，看起來很專業",
        "有滑鼠編輯介面 — 美編也能直接改細節",
    ],
    [
        "只有畫面，沒有 UX — 按鈕點了沒反應",
        "要另外付錢（加 Claude 一個月大概 $36-40）",
    ],
    "畫面最漂亮 + 能細修。但沒有互動。",
    "Paper 專案", "app.paper.design/file/01KPMV...",
)

slide_tool(
    "Stitch AI", STITCH,
    "免費、產出最快、8 種匯出方式。",
    [
        "Google Labs 出的，目前完全免費（每月 350+200 次）",
        "可以直接點文字編輯（其他要用指令）",
        "有 8 種匯出：Figma、.zip、Jules、MCP、預覽連結等",
    ],
    [
        "只能用指令編輯，沒有滑鼠介面",
        "按鈕點了沒反應 — 沒有換頁、沒有彈窗",
    ],
    "最快做草稿的工具。但不能做互動。",
    "Stitch 專案", "stitch.withgoogle.com/projects/...",
)

slide_tool(
    "AI Studio", AISTUDIO,
    "直接做能用的完整網站，免費。",
    [
        "一開始有多種風格可以選",
        "有滑鼠介面改顏色、字型、大小",
        "可以一鍵發布到 Cloud Run，或下載 .zip 自己放 Cloudflare",
    ],
    [
        "畫面比較陽春，不像專業設計工具那麼漂亮",
        "不適合精細的視覺設計",
    ],
    "做「能用的」網站最快。視覺品質就還好。",
    "AI Studio 專案", "ai.studio/apps/57bb988c...",
)

slide_tool(
    "Claude Design", CLAUDE,
    "互動最完整 + 編輯方式最多。已經有 Claude Pro 就免費。",
    [
        "生成前會先問你問題 — 確認你想要什麼",
        "有真的互動：按鈕點下去會跳出彈窗",
        "可以用指令、滑鼠、畫圖、留言、調整滑桿來改",
    ],
    [
        "每週有額度限制，用太多要等下週",
        "只能用瀏覽器，沒有桌面版，斷網就不能用",
    ],
    "互動最強 + 編輯方式最多。Claude Pro 用戶免費。",
    "Claude Design 專案", "claude.ai/design/p/decda9f3...",
)

slide_figma_pivot()
slide_final()

# ============= SAVE =============
out = r"C:\Users\user\Documents\GitHub\TzuchiDigital\uiuxsoftware-evaluation\design-tools-evaluation-zh.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
